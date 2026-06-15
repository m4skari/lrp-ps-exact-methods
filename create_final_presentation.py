"""Build the final English PowerPoint deck for the LRP-PS project."""

from __future__ import annotations

import json
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
RESULTS = ROOT / "results"
OUT = ROOT / "LRP_PS_BnP_vs_BC_MCI_Presentation.pptx"
NOTES = ROOT / "PRESENTATION_NOTES_EN.md"
PLOT_DIR = RESULTS / "07_presentation"

NAVY = RGBColor(22, 42, 68)
TEAL = RGBColor(0, 132, 150)
ORANGE = RGBColor(226, 124, 45)
BLUE = RGBColor(54, 103, 165)
GREEN = RGBColor(44, 140, 92)
RED = RGBColor(188, 70, 70)
GRAY = RGBColor(90, 100, 110)
LIGHT = RGBColor(246, 248, 251)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(35, 42, 50)
MID = RGBColor(210, 219, 229)


def add_text(slide, x, y, w, h, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, bullets, size=17, color=DARK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    for idx, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.name = "Aptos"
        p.font.size = Pt(size - 2 * level)
        p.font.color.rgb = color
        p.space_after = Pt(7 if level == 0 else 3)
    return box


def add_header(slide, title, section=""):
    add_text(slide, 0.55, 0.22, 11.2, 0.52, title, 25, NAVY, True)
    if section:
        add_text(slide, 10.9, 0.33, 1.85, 0.25, section.upper(), 9, TEAL, True, PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.86), Inches(12.2), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def add_footer(slide, n):
    add_text(slide, 0.58, 7.08, 7.2, 0.22, "LRP-PS: Branch-and-Price vs Branch-and-Cut with MCI", 8, GRAY)
    add_text(slide, 12.05, 7.06, 0.7, 0.22, str(n), 8, GRAY, False, PP_ALIGN.RIGHT)


def blank(prs, title, section=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_header(slide, title, section)
    add_footer(slide, len(prs.slides))
    return slide


def add_card(slide, x, y, w, h, title, body, color=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = MID
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = color
    stripe.line.fill.background()
    add_text(slide, x + 0.18, y + 0.12, w - 0.3, 0.32, title, 15, color, True)
    add_text(slide, x + 0.18, y + 0.55, w - 0.3, h - 0.65, body, 12.5, DARK)


def add_picture(slide, path, x, y, w, h):
    path = Path(path)
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    return slide.shapes.add_picture(str(path), Inches(x + (w - pw) / 2), Inches(y + (h - ph) / 2), Inches(pw), Inches(ph))


def render_equation(name, tex, width=10.5, height=0.55, fontsize=23, color="#162a44"):
    PLOT_DIR.mkdir(parents=True, exist_ok=True)
    path = PLOT_DIR / f"eq_{name}.png"
    fig = plt.figure(figsize=(width, height))
    fig.patch.set_alpha(0)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.axis("off")
    ax.text(0.5, 0.5, tex, ha="center", va="center", fontsize=fontsize, color=color)
    fig.savefig(path, dpi=220, transparent=True, bbox_inches="tight", pad_inches=0.04)
    plt.close(fig)
    return path


def add_equation(slide, name, tex, x, y, w, h, fontsize=23, color="#162a44"):
    return add_picture(slide, render_equation(name, tex, w, h, fontsize, color), x, y, w, h)


def add_table(slide, x, y, w, h, dataframe, font_size=9):
    rows, cols = dataframe.shape[0] + 1, dataframe.shape[1]
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for j, col in enumerate(dataframe.columns):
        table.cell(0, j).text = str(col)
    for i, (_, row) in enumerate(dataframe.iterrows(), start=1):
        for j, col in enumerate(dataframe.columns):
            val = row[col]
            if isinstance(val, float):
                table.cell(i, j).text = f"{val:.3f}"
            else:
                table.cell(i, j).text = str(val)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r == 0 else WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                p.font.name = "Aptos"
                p.font.size = Pt(font_size)
                p.font.bold = r == 0
                p.font.color.rgb = WHITE if r == 0 else DARK
    return table


def make_extra_plots():
    PLOT_DIR.mkdir(parents=True, exist_ok=True)
    comp = pd.read_csv(RESULTS / "06_comparison" / "comparison_results.csv")
    svc = pd.read_csv(RESULTS / "06_comparison" / "service_split_results.csv")
    bnp = comp[comp.method == "paper_branch_price"].set_index("instance")
    bc = comp[comp.method == "branch_cut_mci"].set_index("instance")
    names = bnp.index.tolist()
    short = [n.replace("random_", "R-").replace("cluster_", "C-").replace("center_", "ctr-").replace("corner_", "cor-") for n in names]

    fig, ax = plt.subplots(figsize=(10, 4.8))
    x = np.arange(len(names))
    width = 0.36
    ax.bar(x - width / 2, bnp.runtime, width, label="Paper B&P", color="#3667a5")
    ax.bar(x + width / 2, bc.runtime, width, label="B&C + MCI", color="#e27c2d")
    ax.set_yscale("log")
    ax.set_xticks(x, short, rotation=18, ha="right")
    ax.set_ylabel("Runtime (seconds, log scale)")
    ax.set_title("Runtime comparison")
    ax.grid(axis="y", alpha=0.25)
    ax.legend()
    fig.tight_layout()
    runtime = PLOT_DIR / "runtime.png"
    fig.savefig(runtime, dpi=190)
    plt.close(fig)

    fig, ax = plt.subplots(figsize=(10, 4.8))
    ax.plot(short, bnp.objective, marker="o", label="Paper B&P", color="#3667a5")
    ax.plot(short, bc.objective, marker="s", label="B&C + MCI", color="#e27c2d")
    ax.set_ylabel("Objective value")
    ax.set_title("Both exact methods reach the same optimum")
    ax.grid(alpha=0.25)
    ax.legend()
    fig.tight_layout()
    obj = PLOT_DIR / "objectives.png"
    fig.savefig(obj, dpi=190)
    plt.close(fig)

    fig, ax = plt.subplots(figsize=(10, 4.8))
    ax.bar(short, bc.mci_candidates.fillna(0), label="Candidate minimal covers", color="#3667a5")
    ax.bar(short, bc.mci_cuts_added.fillna(0), label="Violated cuts added", color="#e27c2d")
    ax.set_ylabel("Count")
    ax.set_title("MCI generation and separation")
    ax.grid(axis="y", alpha=0.25)
    ax.legend()
    fig.tight_layout()
    mci = PLOT_DIR / "mci_counts.png"
    fig.savefig(mci, dpi=190)
    plt.close(fig)

    paper_svc = svc[svc.method == "paper_branch_price"].set_index("instance").loc[names]
    fig, ax = plt.subplots(figsize=(10, 4.8))
    ax.bar(short, paper_svc.facility_demand, label="Facilities", color="#e27c2d")
    ax.bar(short, paper_svc.small_gv_demand, bottom=paper_svc.facility_demand, label="Small GVs", color="#3667a5")
    ax.set_ylabel("Demand units")
    ax.set_title("Demand split in the optimal solution")
    ax.grid(axis="y", alpha=0.25)
    ax.legend()
    fig.tight_layout()
    split = PLOT_DIR / "demand_split.png"
    fig.savefig(split, dpi=190)
    plt.close(fig)

    return {"runtime": runtime, "obj": obj, "mci": mci, "split": split}


def build_deck():
    plots = make_extra_plots()
    comp = pd.read_csv(RESULTS / "06_comparison" / "comparison_results.csv")
    svc = pd.read_csv(RESULTS / "06_comparison" / "service_split_results.csv")
    mci_bound = pd.read_csv(RESULTS / "06_comparison" / "mci_bound_analysis.csv")
    instances = json.loads((RESULTS / "01_data" / "instances.json").read_text())

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.22), Inches(7.5))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = TEAL
    stripe.line.fill.background()
    add_text(slide, 0.8, 1.15, 11.7, 1.5, "Location-Routing Problem with Pick-up Stations", 35, WHITE, True)
    add_text(slide, 0.82, 2.85, 11.2, 0.7, "Mathematical Model, Branch-and-Price, and Branch-and-Cut with Minimal Cover Inequalities", 20, RGBColor(210, 228, 238))
    add_text(slide, 0.85, 5.25, 10.5, 0.35, "Based on Wang, Zhang, Bell, and Miao (EJOR 2022)", 14, WHITE)
    add_text(slide, 0.85, 5.75, 10.5, 0.35, "Computational implementation and experiments", 14, WHITE)
    add_footer(slide, 1)

    slide = blank(prs, "Roadmap", "overview")
    add_bullets(slide, 1.05, 1.25, 11.4, 5.2, [
        "Problem definition: delivery options, cost components, and assumptions.",
        "Compact mathematical formulation: objective function and all constraint groups.",
        "Paper decomposition: set-partitioning master, feasible columns, and pricing subproblems.",
        "Solution algorithms: paper Branch-and-Price and Branch-and-Cut with MCI.",
        "Data generation and experimental considerations.",
        "Comparison of objective values, runtimes, MCI counts, and service shares.",
        "Code availability and reproducibility.",
    ], 21)

    slide = blank(prs, "Problem Statement: LRP-PS", "problem")
    add_picture(slide, RESULTS / "01_data" / "cluster_corner_n12_p3.png", 0.55, 1.15, 5.9, 5.45)
    add_card(slide, 6.8, 1.1, 5.7, 1.2, "Network", "Depot 0, customer set Vc, and candidate pick-up station set Vl.", TEAL)
    add_card(slide, 6.8, 2.55, 5.7, 1.35, "Delivery rule", "Every customer demand is satisfied exactly once: either by one opened station or by a small green vehicle route.", ORANGE)
    add_card(slide, 6.8, 4.15, 5.7, 1.55, "Large GV assumption", "Large green vehicles do not visit customer demand points. They only replenish opened stations, represented by the effective fixed station cost.", BLUE)
    add_text(slide, 6.9, 6.1, 5.5, 0.4, "Separate strategy: facility service and direct customer routing are modeled separately.", 14, RED, True, PP_ALIGN.CENTER)

    slide = blank(prs, "Notation", "model")
    add_card(slide, 0.65, 1.1, 3.8, 5.2, "Sets", "0: depot\nVc: customers\nVl: candidate pick-up stations\nA: small-GV arcs among depot and customers only", TEAL)
    add_card(slide, 4.75, 1.1, 3.8, 5.2, "Parameters", "q_i: customer demand\nQ_l: station capacity\nQ_e: small-GV capacity\nB: battery range\nr_l: station coverage radius\nc_ij: travel cost\nf_l, a_l: station opening/handling costs\nphi: large-GV route saving factor", ORANGE)
    add_card(slide, 8.85, 1.1, 3.8, 5.2, "Decision variables", "v_l: station opened\nz_il: customer i assigned to station l\nw_i: customer i served by small GV\nx_ij: small GV traverses arc (i,j)\nf_ij: load flow\nb_i: remaining battery", BLUE)

    slide = blank(prs, "Objective Function", "model")
    add_equation(
        slide,
        "objective",
        r"$\min\; \sum_{\ell\in V_l}\bar f_\ell v_\ell"
        r"+\sum_{i\in V_c}\sum_{\ell\in V_l} a_\ell q_i z_{i\ell}"
        r"+c_e\sum_{j\in V_c} x_{0j}"
        r"+\sum_{i\in V_c\cup\{0\}}\sum_{j\in V_c\cup\{0\}:j\neq i} c_{ij}x_{ij}$",
        0.65,
        1.12,
        12.1,
        0.95,
        22,
    )
    add_equation(
        slide,
        "station_fixed_cost",
        r"$\bar f_\ell=f_\ell+\phi(c_{0\ell}+c_{\ell0})$",
        2.1,
        2.18,
        9.0,
        0.55,
        25,
        "#bc4646",
    )
    add_bullets(slide, 1.0, 3.25, 11.4, 2.65, [
        "The first term opens and replenishes pick-up stations.",
        "The second term charges handling cost for packages assigned to stations.",
        "The third term charges the fixed cost of each used small GV.",
        "The fourth term charges the travel cost of small-GV customer routes.",
        "Large-GV travel saving is modeled by phi; if phi = 0.1, only 10% of the station round trip is charged to LRP-PS.",
    ], 18)

    slide = blank(prs, "Assignment and Facility Constraints", "model")
    add_text(slide, 0.9, 1.1, 5.7, 0.6, "(2) Exactly one delivery mode", 18, TEAL, True)
    add_equation(slide, "assign_one_mode", r"$\sum_{\ell\in V_l}z_{i\ell}+w_i=1\qquad \forall i\in V_c$", 0.95, 1.65, 5.45, 0.6, 20)
    add_text(slide, 0.9, 2.75, 5.7, 0.6, "(3) Coverage and station opening", 18, ORANGE, True)
    add_equation(slide, "coverage_opening", r"$d_{i\ell}z_{i\ell}\leq r_\ell v_\ell\qquad \forall i\in V_c,\ell\in V_l$", 0.95, 3.25, 5.45, 0.6, 20)
    add_text(slide, 6.9, 1.1, 5.4, 0.6, "(4) Station capacity", 18, BLUE, True)
    add_equation(slide, "station_capacity", r"$\sum_{i\in V_c}q_i z_{i\ell}\leq Q_\ell v_\ell\qquad \forall \ell\in V_l$", 6.85, 1.65, 5.45, 0.6, 20)
    add_bullets(slide, 6.8, 2.8, 5.5, 2.7, [
        "If v_l = 0, no customer can be assigned to station l.",
        "A station can serve only customers inside its coverage radius.",
        "Total assigned demand cannot exceed station accommodation capacity.",
        "These constraints are the source of minimal cover inequalities.",
    ], 16)

    slide = blank(prs, "Small-GV Routing Constraints", "model")
    add_text(slide, 0.75, 1.05, 5.8, 0.5, "Depot balance and fleet limit", 17, TEAL, True)
    add_equation(slide, "depot_fleet", r"$\sum_j x_{0j}=\sum_j x_{j0},\qquad \sum_jx_{0j}\leq N_e$", 0.75, 1.55, 5.8, 0.68, 18)
    add_text(slide, 6.8, 1.05, 5.8, 0.5, "Customer flow conservation", 17, ORANGE, True)
    add_equation(slide, "customer_flow", r"$\sum_jx_{ij}=\sum_jx_{ji},\qquad \sum_jx_{ij}\leq 1\quad \forall i\in V_c$", 6.8, 1.55, 5.65, 0.68, 17)
    add_text(slide, 0.75, 3.1, 5.8, 0.5, "Load flow and capacity", 17, BLUE, True)
    add_equation(slide, "load_flow", r"$\sum_j f_{ji}-\sum_j f_{ij}=q_iw_i,\qquad f_{ij}\leq Q_ex_{ij}$", 0.65, 3.55, 5.95, 0.75, 17)
    add_text(slide, 6.8, 3.1, 5.8, 0.5, "Battery range", 17, RED, True)
    add_equation(slide, "battery_1", r"$b_i\leq B-hd_{0i}x_{0i}$", 6.95, 3.48, 5.2, 0.38, 17)
    add_equation(slide, "battery_2", r"$b_j\leq b_i-hd_{ij}x_{ij}+B(1-x_{ij})$", 6.95, 3.92, 5.2, 0.5, 17)
    add_text(slide, 1.0, 5.55, 11.3, 0.5, "The x, f, and b variables model only small-GV customer routes. Large GVs are not allowed to visit demand points.", 16, RED, True, PP_ALIGN.CENTER)

    slide = blank(prs, "Paper Decomposition: Main Idea", "paper method")
    add_card(slide, 0.7, 1.25, 3.6, 4.75, "Dantzig-Wolfe Reformulation", "Replace arc-level routing and assignment details by feasible columns.\n\nThe master problem selects columns that cover all customers exactly once.", TEAL)
    add_card(slide, 4.85, 1.25, 3.6, 4.75, "Pattern columns", "A station pattern p is a feasible subset of customers assigned to one station l.\n\nIt must satisfy coverage and station capacity.", ORANGE)
    add_card(slide, 9.0, 1.25, 3.6, 4.75, "Route columns", "A route r is a feasible elementary small-GV route starting and ending at the depot.\n\nIt must satisfy load and battery resources.", BLUE)

    slide = blank(prs, "Master Problem (Set Partitioning)", "paper method")
    add_equation(slide, "master_obj", r"$\min\;\sum_{\ell\in V_l}\bar f_\ell v_\ell+\sum_{p\in P}c_pz_p+\sum_{r\in R}c_rx_r$", 0.8, 1.05, 11.7, 0.65, 22)
    add_equation(slide, "master_cover", r"$\sum_{p\in P}\alpha_{ip}z_p+\sum_{r\in R}\beta_{ir}x_r=1\qquad \forall i\in V_c$", 0.8, 1.95, 11.7, 0.55, 19, "#232a32")
    add_equation(slide, "master_station_link", r"$\sum_{p\in P_\ell}z_p=v_\ell\qquad \forall \ell\in V_l$", 0.8, 2.75, 11.7, 0.5, 19, "#232a32")
    add_equation(slide, "master_fleet", r"$\sum_{r\in R}x_r=H_e,\qquad 0\leq H_e\leq N_e$", 0.8, 3.5, 11.7, 0.5, 19, "#232a32")
    add_bullets(slide, 0.9, 5.05, 11.5, 1.2, [
        "α_ip = 1 if customer i appears in station pattern p.",
        "β_ir = 1 if customer i is visited by small-GV route r.",
        "The master chooses a mutually consistent station-route system.",
    ], 16)

    slide = blank(prs, "Feasible Pattern Columns and PSCSP", "paper method")
    add_text(slide, 0.8, 1.1, 5.7, 0.6, "A feasible station pattern p for station l:", 18, TEAL, True)
    add_bullets(slide, 1.0, 1.8, 5.6, 1.9, [
        "contains only customers inside the coverage radius;",
        "respects the station accommodation capacity;",
        "charges one handling cost term for each assigned demand unit.",
    ], 16)
    add_text(slide, 6.85, 1.1, 5.55, 0.6, "Pricing subproblem PSCSP-l:", 18, ORANGE, True)
    add_equation(slide, "pscsp_obj", r"$\min\;\sum_i(a_\ell q_i-\mu_i)z_i^\ell-\tau_\ell$", 6.75, 1.72, 5.75, 0.45, 17)
    add_equation(slide, "pscsp_cov", r"$d_{i\ell}z_i^\ell\leq r_\ell,\quad \sum_iq_iz_i^\ell\leq Q_\ell,\quad z_i^\ell\in\{0,1\}$", 6.7, 2.32, 5.8, 0.55, 15)
    add_equation(
        slide,
        "pscsp_knapsack",
        r"$\max\;\sum_i(\mu_i-a_\ell q_i)z_i^\ell+\tau_\ell\quad\Rightarrow\quad\mathrm{0\!-\!1\ knapsack}$",
        0.75,
        4.45,
        11.8,
        0.7,
        18,
        "#bc4646",
    )

    slide = blank(prs, "Feasible Route Columns and SPPBDRC", "paper method")
    add_bullets(slide, 0.8, 1.15, 5.8, 2.7, [
        "A feasible route r starts at the depot, visits a subset of customers, and returns to the depot.",
        "The route is elementary: no customer is repeated.",
        "Total demand on the route is at most Q_e.",
        "Total energy consumption is within battery range B.",
        "The route cost is c_r = c_e + travel distance.",
    ], 16)
    add_text(slide, 6.8, 1.2, 5.6, 0.35, "Reduced arc costs:", 17, NAVY, True, PP_ALIGN.CENTER)
    add_equation(slide, "route_reduced_1", r"$\bar c_{ij}=c_{ij}-\mu_i\quad (i\in V_c)$", 6.9, 1.75, 5.35, 0.45, 17)
    add_equation(slide, "route_reduced_2", r"$\bar c_{0j}=c_e+c_{0j}-\eta$", 6.9, 2.25, 5.35, 0.45, 17)
    add_text(slide, 6.8, 3.25, 5.6, 0.9, "Pricing: find a negative reduced-cost elementary shortest path with load and battery resources.", 18, ORANGE, True, PP_ALIGN.CENTER)
    add_text(slide, 6.8, 4.75, 5.6, 0.7, "Known problem class: ESPPRC / resource-constrained shortest path.", 18, RED, True, PP_ALIGN.CENTER)

    slide = blank(prs, "Column Generation and Branching", "paper method")
    add_bullets(slide, 0.85, 1.2, 5.8, 4.8, [
        "Start from a restricted master problem.",
        "Solve the LP relaxation and read dual variables μ, τ, η.",
        "Solve PSCSP and SPPBDRC pricing problems.",
        "Add columns with negative reduced cost.",
        "Repeat until no improving column exists.",
        "Branch when the master solution is fractional.",
    ], 17)
    add_picture(slide, RESULTS / "04_paper_method" / "cluster_corner_n12_p3_convergence.png", 6.8, 1.25, 5.7, 4.5)
    add_text(slide, 6.9, 6.0, 5.5, 0.4, "The implemented benchmark completes route columns exactly for small instances and then solves the integer master.", 13, GRAY, False, PP_ALIGN.CENTER)

    slide = blank(prs, "Branch-and-Cut with Minimal Cover Inequalities", "B&C + MCI")
    add_equation(slide, "cover_def", r"$C\subseteq E_\ell\ \mathrm{is\ a\ cover\ if}\ \sum_{i\in C}q_i>Q_\ell$", 0.8, 1.08, 11.8, 0.6, 21)
    add_equation(slide, "minimal_cover_def", r"$C\ \mathrm{is\ minimal\ if}\ \sum_{i\in C\setminus\{k\}}q_i\leq Q_\ell\quad \forall k\in C$", 0.8, 1.95, 11.8, 0.65, 19, "#232a32")
    add_equation(slide, "mci_main", r"$\sum_{i\in C}z_{i\ell}\leq (|C|-1)v_\ell$", 1.4, 3.12, 10.4, 0.75, 28, "#bc4646")
    add_bullets(slide, 1.0, 4.55, 11.1, 1.35, [
        "If v_l = 0, no assignment to station l is possible.",
        "If v_l = 1, at most |C|-1 customers from an over-capacity cover can be assigned.",
        "The inequality is valid for all integer feasible solutions and can cut fractional node relaxations.",
    ], 16)

    slide = blank(prs, "B&C + MCI Algorithm", "B&C + MCI")
    add_bullets(slide, 0.85, 1.1, 11.8, 5.3, [
        "1. Build the compact MIP with exact preprocessing.",
        "2. Enumerate minimal covers for each station-capacity constraint.",
        "3. Run Gurobi Branch-and-Bound on the integer compact model.",
        "4. At a fractional MIP node, evaluate each candidate cover:",
        ("if sum z*_il > (|C|-1)v*_l, add the user cut sum z_il <= (|C|-1)v_l.", 1),
        "5. Continue branching and cutting until the global MIP gap is zero.",
        "The B&C model is exact because MCI cuts are globally valid.",
    ], 18)

    slide = blank(prs, "Why MCI Counts Differ Across Instances", "B&C + MCI")
    add_picture(slide, RESULTS / "06_comparison" / "mci_cut_activity.png", 0.55, 1.2, 6.2, 4.8)
    add_bullets(slide, 7.05, 1.25, 5.45, 4.8, [
        "Many MCIs appear when several customers are eligible for the same station and their combined demand exceeds Q_l.",
        "Few or zero MCIs appear when station neighborhoods are small or only one/two customers are eligible.",
        "A candidate MCI is added only if it is violated at a fractional node.",
        "Therefore, candidate count depends on geometry and demand/capacity, while added-cut count depends on the LP relaxation path.",
        "The clustered instance has many nearby customers around stations, so it generated 56 candidates and 24 active cuts.",
    ], 15)

    slide = blank(prs, "Data Generation and Modeling Considerations", "experiments")
    inst_df = pd.DataFrame(instances)
    inst_df = inst_df[["name", "customers", "stations", "total_demand", "vehicle_capacity", "station_capacity", "coverage_radius", "max_small_gvs"]]
    inst_df.columns = ["Instance", "n", "p", "Demand", "Qe", "Ql", "r", "Ne"]
    add_table(slide, 0.45, 1.1, 12.4, 2.3, inst_df, 8)
    add_bullets(slide, 0.75, 3.8, 11.8, 2.3, [
        "Customer coordinates are generated with three spatial structures: central-random, corner-random, and corner-clustered.",
        "Stations are placed near demand centers, then coverage is tightened: r_l = clip(0.14B, 12, 22).",
        "The number of small GVs is explicitly capped: N_e = ceil(total demand / Q_e) + 1.",
        "Large-GV station replenishment is represented through f_bar_l and the phi travel-saving factor.",
    ], 15)

    slide = blank(prs, "Optimal Demand Split and Large-GV Saving", "experiments")
    add_picture(slide, PLOT_DIR / "demand_split.png", 0.55, 1.2, 6.4, 4.9)
    paper_svc = svc[svc.method == "paper_branch_price"].copy()
    paper_svc["Facility %"] = 100 * paper_svc["facility_share"]
    paper_svc["Small GV %"] = 100 * paper_svc["small_gv_share"]
    paper_svc["Saving"] = paper_svc["large_gv_travel_saving"]
    tab = paper_svc[["instance", "Facility %", "Small GV %", "small_gv_routes", "max_small_gvs", "Saving"]]
    tab.columns = ["Instance", "Fac.%", "Small.%", "Routes", "Ne", "LGV saving"]
    add_table(slide, 7.1, 1.35, 5.75, 4.5, tab, 7.2)
    add_text(slide, 7.25, 6.1, 5.3, 0.35, "The fleet cap is respected in all optimal solutions.", 13, GREEN, True, PP_ALIGN.CENTER)

    slide = blank(prs, "Objective Values: Two Exact Methods", "results")
    add_picture(slide, PLOT_DIR / "objectives.png", 0.55, 1.1, 6.4, 4.9)
    obj = comp.pivot(index="instance", columns="method", values="objective").reset_index()
    obj["Abs. diff"] = (obj["paper_branch_price"] - obj["branch_cut_mci"]).abs()
    obj = obj[["instance", "paper_branch_price", "branch_cut_mci", "Abs. diff"]]
    obj.columns = ["Instance", "Paper B&P", "B&C + MCI", "|Diff|"]
    add_table(slide, 7.1, 1.45, 5.75, 3.7, obj, 8)
    add_text(slide, 7.25, 5.65, 5.25, 0.5, "Both algorithms reach the same certified optimum on every instance.", 15, GREEN, True, PP_ALIGN.CENTER)

    slide = blank(prs, "Runtime Comparison", "results")
    add_picture(slide, PLOT_DIR / "runtime.png", 0.55, 1.1, 7.2, 5.2)
    rt = comp.pivot(index="instance", columns="method", values="runtime").reset_index()
    rt["B&C / B&P"] = rt["branch_cut_mci"] / rt["paper_branch_price"]
    rt = rt[["instance", "paper_branch_price", "branch_cut_mci", "B&C / B&P"]]
    rt.columns = ["Instance", "B&P sec", "B&C sec", "Ratio"]
    add_table(slide, 8.0, 1.5, 4.9, 3.6, rt, 8)
    add_text(slide, 8.15, 5.55, 4.5, 0.55, "The paper decomposition is faster because it optimizes over feasible patterns/routes directly instead of branching on many arc variables.", 13, RED, True, PP_ALIGN.CENTER)

    slide = blank(prs, "Root LP Effect of MCI", "results")
    add_table(slide, 0.55, 1.1, 12.25, 3.0, mci_bound[["instance", "mci_count", "raw_gap_pct", "mci_gap_pct", "bound_improvement_pct"]].rename(columns={"instance": "Instance", "mci_count": "MCI", "raw_gap_pct": "Raw gap %", "mci_gap_pct": "MCI gap %", "bound_improvement_pct": "Improvement %"}), 8)
    add_bullets(slide, 0.9, 4.65, 11.5, 1.55, [
        "MCI improved the root LP bound only for the clustered instance.",
        "This does not make MCI invalid; it means the violated cover structure was absent or inactive at the root in other instances.",
        "MCI is more useful inside Branch-and-Cut when fractional node solutions violate station capacity-cover structure.",
    ], 17)

    slide = blank(prs, "Conclusions", "conclusion")
    add_bullets(slide, 0.9, 1.1, 11.5, 4.7, [
        "The corrected model keeps large GVs away from demand points; they only replenish opened facilities.",
        "The paper's Branch-and-Price is the best-performing exact method in these experiments.",
        "Branch-and-Cut with MCI is exact and reaches the same optimal objectives, but it is slower on harder instances.",
        "MCI activity is highly instance-dependent: many nearby customers plus tight station capacity create many covers.",
        "Demand split, fleet cap compliance, and large-GV travel saving are explicitly reported for each instance.",
        "The code is prepared for GitHub publication with reproducible data generation, experiments, plots, and tests.",
    ], 18)
    add_text(slide, 1.2, 6.15, 10.8, 0.5, "GitHub repository link: to be added after remote repository creation and push.", 17, TEAL, True, PP_ALIGN.CENTER)

    prs.save(OUT)
    write_notes()
    return OUT


def write_notes():
    NOTES.write_text(
        """# Speaker Notes: LRP-PS Presentation

The slides are intentionally math-heavy. The main points to emphasize are:

1. Large GVs do not visit customers. They only replenish opened pick-up stations.
2. Customer demand is satisfied either by a station or by a small GV route, never both.
3. The objective includes station opening, discounted large-GV station round trip, station handling, small-GV fixed cost, and small-GV travel.
4. Constraints (2)-(4) are assignment, coverage/opening, and station capacity constraints.
5. Routing constraints are only for small GVs: depot balance, fleet cap, customer flow, load flow, and battery propagation.
6. In the paper decomposition, station patterns become knapsack-type columns and small-GV routes become ESPPRC-type columns.
7. The master problem is a set-partitioning model over feasible columns.
8. MCI cuts are derived from minimal covers of station-capacity constraints.
9. B&C + MCI is exact because these cuts are globally valid.
10. MCI count is large when many customers are near the same station and station capacity is tight.
11. MCI count is small when neighborhoods are sparse, station capacity is loose, or root/node LP solutions do not violate the covers.
12. In this benchmark, B&P is faster than B&C + MCI, but both certify the same optimum.
13. GitHub push was attempted separately; if no credential is available, add the final repository URL manually after upload.
""",
        encoding="utf-8",
    )


if __name__ == "__main__":
    print(build_deck())
